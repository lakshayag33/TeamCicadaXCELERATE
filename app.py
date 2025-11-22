from flask import Flask, render_template, request, redirect, url_for, session, send_file
from flask_sqlalchemy import SQLAlchemy
import os
from werkzeug.utils import secure_filename
from rag_engine import RAGEngine
from llm_client import GraniteClient
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from io import BytesIO
from textwrap import wrap


# Import libraries for non-PDF document handling
from docx import Document
from pptx import Presentation

# ====================================================
# CONFIGURATION
# ====================================================
app = Flask(__name__)
app.secret_key = 'studymate_secret_key_2025'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///history.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'pptx', 'txt'}

db = SQLAlchemy(app)

# ====================================================
# DATABASE MODEL
# ====================================================
class QuestionHistory(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    question = db.Column(db.Text, nullable=False)
    answer = db.Column(db.Text, nullable=False)

    def __repr__(self):
        return f"<QuestionHistory {self.id}>"

# ====================================================
# INITIALIZATION
# ====================================================
rag_engine = RAGEngine(chunk_size=250, chunk_overlap=120, debug=True)
llm_client = GraniteClient(device='cuda')  # GPU if available

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
with app.app_context():
    db.create_all()

# ====================================================
# HELPER FUNCTIONS
# ====================================================
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def extract_text_from_docx(file_path):
    """Extract text from Word document"""
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        print(f"[DOCX ERROR] {e}")
        return ""

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentation"""
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        print(f"[PPTX ERROR] {e}")
        return ""

def extract_text_from_txt(file_path):
    """Extract text from plain text file"""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"[TXT ERROR] {e}")
        return ""
    
    from textwrap import wrap  # add at the top with other imports if not already

def draw_wrapped_text(pdf, text, x, y, max_width, line_height=13):
    """
    Draw text on the PDF with automatic word wrapping and page breaks.
    """
    wrapped_lines = wrap(text, width=max_width)
    for line in wrapped_lines:
        if y < 80:  # Start a new page if bottom reached
            pdf.showPage()
            pdf.setFont("Helvetica", 10)
            y = letter[1] - 60
        pdf.drawString(x, y, line)
        y -= line_height
    return y


# ====================================================
# ROUTES
# ====================================================
@app.route('/')
def index():
    session.clear()
    return render_template('index.html')


@app.route('/process', methods=['POST'])
def process():
    print("\n" + "=" * 90)
    print("[APP] Starting file processing...")
    print("=" * 90)

    if 'pdfs' not in request.files:
        print("[APP] ERROR: No 'pdfs' in request.files")
        return redirect(url_for('index'))

    files = request.files.getlist('pdfs')
    print(f"[APP] Received {len(files)} file(s)")

    if not files or files[0].filename == '':
        print("[APP] ERROR: Empty file list or first file has no name")
        return redirect(url_for('index'))

    text_files = []
    for file in files:
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            ext = filename.rsplit('.', 1)[1].lower()
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            print(f"[APP] Saved: {filepath}")

            # Extract text based on file type
            extracted_text = ""
            if ext == "pdf":
                text_files.append(filepath)
                print(f"[APP] Using PDF directly: {filepath}")
                continue
            elif ext == "docx":
                extracted_text = extract_text_from_docx(filepath)
            elif ext == "pptx":
                extracted_text = extract_text_from_pptx(filepath)
            elif ext == "txt":
                extracted_text = extract_text_from_txt(filepath)
            else:
                print(f"[APP] Skipping unsupported file type: {ext}")
                continue

            # Save extracted content as a temporary .txt file
            temp_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(filename)[0]}_converted.txt")
            with open(temp_path, "w", encoding="utf-8") as f:
                f.write(extracted_text)
            text_files.append(temp_path)
            os.remove(filepath)  # remove original non-pdf file
            print(f"[APP] Converted and saved as: {temp_path}")

    if not text_files:
        print("[APP] ERROR: No valid or convertible files found.")
        return redirect(url_for('index'))

    try:
        print(f"[APP] Processing {len(text_files)} file(s) with RAG...")
        rag_engine.process_pdfs(text_files)
        session['processed'] = True
        session['num_chunks'] = len(rag_engine.chunks)
        print(f"[APP] ✅ RAG ready with {session['num_chunks']} chunks")

        # Clean up temp files
        for path in text_files:
            if os.path.exists(path):
                os.remove(path)
                print(f"[APP] Removed temp: {path}")

        print("[APP] Processing complete.\n")
        return render_template('ask.html', num_chunks=session['num_chunks'])
    except Exception as e:
        import traceback
        print("[APP] ERROR during processing:", e)
        print(traceback.format_exc())
        return f"Error processing files: {str(e)}", 500


@app.route('/ask', methods=['GET'])
def ask_page():
    if not session.get('processed'):
        return redirect(url_for('index'))
    return render_template('ask.html', num_chunks=session.get('num_chunks', 0))


@app.route('/answer', methods=['POST'])
def answer():
    if not session.get('processed'):
        return redirect(url_for('index'))

    question = request.form.get('question', '').strip()
    if not question:
        return redirect(url_for('ask_page'))

    try:
        print("\n" + "=" * 90)
        print("[APP] Question:", question)
        print("=" * 90)

        # Retrieve context
        context, ranked = rag_engine.build_context(question, top_k=5, max_chars=3000)
        print(f"[APP] Context length: {len(context)}")

        if not context.strip():
            answer_text = "Information not found in the document."
            sources = []
        else:
            print("[APP] Calling Granite to generate answer...")
            answer_text = llm_client.generate_answer(question, context, max_new_tokens=150)
            print("[APP] Answer Generated:", answer_text[:150], "...")

        # ✅ Save to database
        new_entry = QuestionHistory(question=question, answer=answer_text)
        db.session.add(new_entry)
        db.session.commit()
        print("[DB] Saved question to history table.")

        # Prepare sources
        sources = []
        for r in ranked:
            meta = r.get('meta', {})
            src_name = os.path.basename(meta.get('source', ''))
            pages = f"{meta.get('page_start', '?')}-{meta.get('page_end', '?')}"
            sources.append({
                'text': (r['text'][:300] + '...') if len(r['text']) > 300 else r['text'],
                'score': f"{r.get('combined', r.get('score', 0.0)):.3f}",
                'meta': f"{src_name} (pp. {pages})"
            })

        return render_template('answer.html', question=question, answer=answer_text, sources=sources)

    except Exception as e:
        import traceback
        print("[APP] CRITICAL:", e)
        print(traceback.format_exc())
        return f"Critical error: {str(e)}", 500


@app.route('/new_question')
def new_question():
    return redirect(url_for('ask_page'))


@app.route('/reset')
def reset():
    print("[APP] Resetting RAG engine...")
    rag_engine.reset()
    session.clear()
    return redirect(url_for('index'))


@app.route('/history')
def history():
    records = QuestionHistory.query.order_by(QuestionHistory.id.desc()).all()
    return render_template('history.html', records=records)


@app.route('/download_history')
def download_history():
    records = QuestionHistory.query.order_by(QuestionHistory.id.asc()).all()

    if not records:
        return "No history available to download.", 404

    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    y = height - 60

    pdf.setFont("Helvetica-Bold", 16)
    pdf.drawString(200, y, "📚 StudyMate - Question History")
    y -= 40

    pdf.setFont("Helvetica", 11)
    for idx, record in enumerate(records, start=1):
        q_text = f"Q{idx}: {record.question}"
        a_text = f"A{idx}: {record.answer}"

        pdf.setFont("Helvetica-Bold", 11)
        pdf.drawString(50, y, q_text)
        y -= 15

        pdf.setFont("Helvetica", 10)
        y = draw_wrapped_text(pdf, q_text, 60, y, max_width=85, line_height=18)

        y -= 20
        if y < 80:
            pdf.showPage()
            pdf.setFont("Helvetica", 11)
            y = height - 60

    pdf.save()
    buffer.seek(0)
    return send_file(buffer, as_attachment=True, download_name="StudyMate_History.pdf", mimetype="application/pdf")

# ====================================================
# QUIZ GENERATION FEATURE
# ====================================================
@app.route('/generate_quiz')
def generate_quiz():
    if not session.get('processed'):
        return redirect(url_for('index'))

    try:
        context, _ = rag_engine.build_context("Generate quiz questions", top_k=10, max_chars=4000)
        print(f"[QUIZ] Using {len(context)} characters of context for quiz generation.")

        prompt = (
            "You are an expert exam question designer. Based on the study material below, "
            "generate 5 **high-quality questions** that test understanding, reasoning, and application. "
            "Avoid superficial questions. Make them concise, numbered 1 to 5, and directly related to the text.\n\n"
            f"Study Material:\n{context}"
        )

        quiz_text = llm_client.generate_answer(prompt, context, max_new_tokens=300)
        questions = [q.strip(" .") for q in quiz_text.split("\n") if q.strip()]
        questions = [q for q in questions if len(q) > 5][:5]
        session['quiz'] = questions

        print(f"[QUIZ] Generated {len(questions)} questions.")
        return render_template('quiz.html', questions=questions)

    except Exception as e:
        import traceback
        print("[QUIZ] ERROR during quiz generation:", e)
        print(traceback.format_exc())
        return f"Quiz generation failed: {e}", 500


@app.route('/download_quiz')
def download_quiz():
    quiz = session.get('quiz', [])
    if not quiz:
        return "No quiz available to download.", 404

    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    y = height - 60

    pdf.setFont("Helvetica-Bold", 16)
    pdf.drawString(200, y, "📘 StudyMate - Quiz")
    y -= 40

    pdf.setFont("Helvetica", 12)
    for idx, question in enumerate(quiz, start=1):
        q_text = f"Q{idx}: {question}"
        y = draw_wrapped_text(pdf, q_text, 60, y, max_width=85, line_height=18)

        y -= 10

    pdf.save()
    buffer.seek(0)
    return send_file(buffer, as_attachment=True, download_name="StudyMate_Quiz.pdf", mimetype="application/pdf")

# ====================================================
# ENTRY POINT
# ====================================================
if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
